import tkinter as tk
from tkinter import filedialog, messagebox, Listbox
from docx import Document
import pandas as pd
from pptx import Presentation
import os

def merge_word_documents(documents, output_file):
    """Função para mesclar documentos DOCX."""
    merged_document = Document()
    for document in documents:
        sub_document = Document(document)
        for element in sub_document.element.body:
            merged_document.element.body.append(element)
    merged_document.save(output_file)

def merge_excel_files(documents, output_file):
    """Função para mesclar arquivos XLSX."""
    combined_data = pd.concat((pd.read_excel(doc) for doc in documents), ignore_index=True)
    combined_data.to_excel(output_file, index=False)

def merge_pptx_files(documents, output_file):
    """Função para mesclar arquivos PPTX."""
    merged_presentation = Presentation()
    for document in documents:
        presentation = Presentation(document)
        for slide in presentation.slides:
            new_slide = merged_presentation.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    new_shape = new_slide.shapes[-1]
                    new_shape.text = shape.text
    merged_presentation.save(output_file)

def merge_documents(documents, output_file):
    """Mescla documentos com base no tipo."""
    file_type = os.path.splitext(documents[0])[1][1:]  # Obtém a extensão do primeiro documento
    if file_type == 'docx':
        merge_word_documents(documents, output_file)
    elif file_type == 'xlsx':
        merge_excel_files(documents, output_file)
    elif file_type == 'pptx':
        merge_pptx_files(documents, output_file)
    else:
        messagebox.showerror("Erro", "Tipo de arquivo não suportado!")
        return

    messagebox.showinfo("Sucesso", f"Documentos mesclados com sucesso em:\n{output_file}")

def add_files():
    """Abre uma janela de seleção de arquivos e adiciona à lista."""
    filetypes = [("Documentos Word", "*.docx"), ("Planilhas Excel", "*.xlsx"), ("Apresentações PowerPoint", "*.pptx")]
    files = filedialog.askopenfilenames(title="Selecione arquivos", filetypes=filetypes)
    for file in files:
        if file not in file_list.get(0, tk.END):
            file_list.insert(tk.END, file)

def drag_and_drop(event):
    """Adiciona arquivos arrastados e soltos à lista."""
    files = root.tk.splitlist(event.data)
    for file in files:
        if (file.endswith(".docx") or file.endswith(".xlsx") or file.endswith(".pptx")) and file not in file_list.get(0, tk.END):
            file_list.insert(tk.END, file)

def merge():
    """Mescla os arquivos selecionados."""
    files = file_list.get(0, tk.END)
    if not files:
        messagebox.showwarning("Aviso", "Nenhum arquivo selecionado!")
        return

    # Determine o tipo de arquivo a partir da extensão do primeiro arquivo na lista
    file_type = os.path.splitext(files[0])[1][1:]  # Obtém a extensão do primeiro arquivo
    if file_type not in ['docx', 'xlsx', 'pptx']:
        messagebox.showerror("Erro", "Tipo de documento inválido!")
        return

    output_file = filedialog.asksaveasfilename(
        defaultextension=f".{file_type}",
        filetypes=[(f"Documentos {file_type.upper()}", f"*.{file_type}")],
        title="Salvar documento mesclado como"
    )
    if output_file:
        merge_documents(files, output_file)

def clear_list():
    """Limpa todos os arquivos da lista."""
    file_list.delete(0, tk.END)

# Configuração da janela principal
root = tk.Tk()
root.title("Mesclar Documentos")
root.geometry("600x400")
root.config(bg="#f0f6fc")

# Cabeçalho
title_label = tk.Label(
    root, text="Mesclar Documentos", font=("Segoe UI", 18, "bold"), bg="#f0f6fc", fg="#2d6187"
)
title_label.pack(pady=10)

# Lista para exibir arquivos
file_list = Listbox(root, width=80, height=10, selectmode=tk.MULTIPLE, bg="white", fg="#2d6187")
file_list.pack(padx=10, pady=10)

# Botões de ação
button_frame = tk.Frame(root, bg="#f0f6fc")
button_frame.pack(pady=10)

add_button = tk.Button(button_frame, text="Adicionar Arquivos", command=add_files, bg="#2d6187", fg="white")
add_button.grid(row=0, column=0, padx=5)

merge_button = tk.Button(button_frame, text="Mesclar Documentos", command=merge, bg="#2d6187", fg="white")
merge_button.grid(row=0, column=1, padx=5)

clear_button = tk.Button(button_frame, text="Limpar Lista", command=clear_list, bg="#2d6187", fg="white")
clear_button.grid(row=0, column=2, padx=5)

# Suporte para drag-and-drop (Windows e Mac)
try:
    root.tk.call('tkdnd::drag_source', 'register', root._w)
    root.drop_target_register('DND_Files')
    root.dnd_bind('<<Drop>>', drag_and_drop)
except tk.TclError:
    pass  # Drag-and-drop não disponível no sistema

# Executa a aplicação
root.mainloop()
